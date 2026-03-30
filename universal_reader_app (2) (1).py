#!/usr/bin/env python3
"""
DocShield — Universal Document Reader  v4.0
Fixed: Network/Axios errors, ValueError, PPTX render, Duplicate columns,
       file_uploader reset bug, Arrow serialization, session state corruption
"""

import streamlit as st
import os, json, hashlib, zipfile, tempfile, io, time, traceback
from datetime import datetime
from pathlib import Path

st.set_page_config(
    page_title="DocShield — Universal Reader",
    page_icon="🛡️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ══════════════════════════════════════════════════════════════
#  MASTER CSS
# ══════════════════════════════════════════════════════════════
st.markdown(r"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Orbitron:wght@400;700;900&family=Exo+2:ital,wght@0,300;0,400;0,500;0,600;0,700;1,400&family=Share+Tech+Mono&display=swap');

:root{
  --c:  #00e5ff;  --p:#9333ea;  --g:#00ff88;
  --y:  #ffd60a;  --r:#ff3b5c;  --o:#ff6b35;
  --tx: #e0f0ff;  --mu:#4a6b8a;
  --gl: rgba(8,20,40,.78); --gl2:rgba(10,25,50,.60);
  --bd: rgba(0,200,255,.13);
}
html,body{margin:0;padding:0;}
.stApp{
  background:
    radial-gradient(ellipse 80% 50% at 15% 18%,rgba(0,80,200,.09) 0%,transparent 55%),
    radial-gradient(ellipse 60% 40% at 85% 82%,rgba(147,51,234,.11) 0%,transparent 55%),
    linear-gradient(180deg,#020408 0%,#050b14 60%,#020810 100%) !important;
  min-height:100vh;
  font-family:'Exo 2',sans-serif !important;
  color:var(--tx) !important;
}
.stApp::before{
  content:'';position:fixed;inset:0;
  background-image:
    linear-gradient(rgba(0,229,255,.022) 1px,transparent 1px),
    linear-gradient(90deg,rgba(0,229,255,.022) 1px,transparent 1px);
  background-size:55px 55px;
  animation:gs 28s linear infinite;
  pointer-events:none;z-index:0;
}
@keyframes gs{0%{background-position:0 0}100%{background-position:55px 55px}}

[data-testid="stSidebar"]{
  background:linear-gradient(180deg,rgba(2,8,20,.98) 0%,rgba(4,12,28,.98) 100%) !important;
  border-right:1px solid rgba(0,229,255,.09) !important;
}
[data-testid="stSidebar"]::before{
  content:'';position:absolute;top:0;left:0;right:0;height:2px;
  background:linear-gradient(90deg,transparent,var(--c),var(--p),transparent);
}
[data-testid="stSidebar"] *{color:var(--tx) !important;}
[data-testid="stSidebar"] .stSlider>div>div>div>div{
  background:linear-gradient(90deg,var(--c),var(--p)) !important;
}
[data-testid="stSidebar"] [role="slider"]{
  background:var(--c) !important;
  box-shadow:0 0 12px rgba(0,229,255,.5) !important;
}
[role="switch"][aria-checked="true"]{
  background-color:var(--c) !important;
  box-shadow:0 0 10px rgba(0,229,255,.45) !important;
}
[role="switch"][aria-checked="false"]{background-color:rgba(255,59,92,.35) !important;}

#MainMenu,footer,header{visibility:hidden;}
.block-container{padding-top:.55rem !important;max-width:1440px;position:relative;z-index:1;}

/* ── HERO ─────────────────────────────────── */
.hero{
  background:linear-gradient(135deg,rgba(0,10,30,.96) 0%,rgba(5,0,40,.96) 45%,rgba(0,20,50,.96) 100%);
  border:1px solid rgba(0,229,255,.20);border-radius:20px;
  padding:1.7rem 2.4rem;margin-bottom:1.1rem;
  position:relative;overflow:hidden;
  box-shadow:0 0 55px rgba(0,229,255,.04),0 0 110px rgba(147,51,234,.04),inset 0 1px 0 rgba(255,255,255,.04);
}
.hero::before{
  content:'';position:absolute;top:0;left:-100%;width:300%;height:2px;
  background:linear-gradient(90deg,transparent 0%,var(--c) 28%,var(--p) 52%,var(--o) 76%,transparent 100%);
  animation:sl 4.2s linear infinite;
}
@keyframes sl{0%{left:-100%}100%{left:100%}}
.hero::after{
  content:'';position:absolute;inset:0;pointer-events:none;
  background:
    radial-gradient(ellipse 55% 70% at 0% 50%,rgba(0,229,255,.06) 0%,transparent 60%),
    radial-gradient(ellipse 35% 60% at 100% 50%,rgba(147,51,234,.08) 0%,transparent 60%);
}
.hero-title{
  font-family:'Orbitron',monospace;font-size:2.15rem;font-weight:900;
  background:linear-gradient(90deg,#00e5ff 0%,#7c3aed 35%,#ec4899 65%,#00e5ff 100%);
  background-size:200%;-webkit-background-clip:text;-webkit-text-fill-color:transparent;
  background-clip:text;animation:sh 4s linear infinite;margin:0;letter-spacing:2px;
  position:relative;z-index:1;
}
@keyframes sh{0%{background-position:0%}100%{background-position:200%}}
.hero-sub{
  font-family:'Share Tech Mono',monospace;font-size:.70rem;
  color:rgba(0,229,255,.52);margin:.4rem 0 0 0;
  letter-spacing:.20em;text-transform:uppercase;position:relative;z-index:1;
}

/* ── STAT CARDS ───────────────────────────── */
.sc{border-radius:14px;padding:1.05rem 1.25rem;margin-bottom:.65rem;
    position:relative;overflow:hidden;backdrop-filter:blur(12px);
    transition:transform .2s,box-shadow .2s;}
.sc:hover{transform:translateY(-3px);}
.sc::before{content:'';position:absolute;top:0;left:0;right:0;height:2px;border-radius:14px 14px 0 0;}
.sc-c{background:rgba(0,229,255,.07);border:1px solid rgba(0,229,255,.22);box-shadow:0 4px 26px rgba(0,229,255,.07);}
.sc-c::before{background:linear-gradient(90deg,transparent,var(--c),transparent);}
.sc-r{background:rgba(255,59,92,.08);border:1px solid rgba(255,59,92,.26);box-shadow:0 4px 26px rgba(255,59,92,.08);}
.sc-r::before{background:linear-gradient(90deg,transparent,var(--r),transparent);}
.sc-y{background:rgba(255,214,10,.07);border:1px solid rgba(255,214,10,.23);box-shadow:0 4px 26px rgba(255,214,10,.07);}
.sc-y::before{background:linear-gradient(90deg,transparent,var(--y),transparent);}
.sc-g{background:rgba(0,255,136,.07);border:1px solid rgba(0,255,136,.21);box-shadow:0 4px 26px rgba(0,255,136,.07);}
.sc-g::before{background:linear-gradient(90deg,transparent,var(--g),transparent);}
.sc-lbl{font-family:'Share Tech Mono',monospace;font-size:.62rem;color:var(--mu);text-transform:uppercase;letter-spacing:.20em;margin-bottom:.38rem;}
.sc-val{font-family:'Orbitron',monospace;font-size:1.9rem;font-weight:700;line-height:1;}

/* ── INFO CARDS ───────────────────────────── */
.ic{background:var(--gl2);border:1px solid var(--bd);border-radius:11px;
    padding:.85rem 1.15rem;margin-bottom:.6rem;backdrop-filter:blur(10px);
    transition:border-color .2s,box-shadow .2s;}
.ic:hover{border-color:rgba(0,229,255,.27);box-shadow:0 0 16px rgba(0,229,255,.05);}
.ic-lbl{font-family:'Share Tech Mono',monospace;font-size:.60rem;color:var(--mu);
        text-transform:uppercase;letter-spacing:.18em;margin-bottom:.28rem;}
.ic-val{font-family:'Share Tech Mono',monospace;font-size:.86rem;color:var(--tx);word-break:break-all;line-height:1.5;}

/* ── VERDICTS ─────────────────────────────── */
.v-safe  {background:linear-gradient(135deg,rgba(0,255,136,.08),rgba(0,200,100,.04));border:1px solid rgba(0,255,136,.36);border-left:4px solid var(--g);border-radius:12px;padding:1.05rem 1.35rem;color:#00ff88;font-family:'Exo 2',sans-serif;box-shadow:0 0 26px rgba(0,255,136,.07);}
.v-warn  {background:linear-gradient(135deg,rgba(255,214,10,.08),rgba(255,160,0,.04)); border:1px solid rgba(255,214,10,.36);border-left:4px solid var(--y);border-radius:12px;padding:1.05rem 1.35rem;color:#ffd60a;font-family:'Exo 2',sans-serif;box-shadow:0 0 26px rgba(255,214,10,.07);}
.v-danger{background:linear-gradient(135deg,rgba(255,59,92,.11),rgba(200,0,50,.07));  border:1px solid rgba(255,59,92,.46);border-left:4px solid var(--r);border-radius:12px;padding:1.05rem 1.35rem;color:#ff3b5c;font-family:'Exo 2',sans-serif;animation:dp 2.3s ease-in-out infinite;}
.v-info  {background:linear-gradient(135deg,rgba(0,229,255,.07),rgba(0,150,255,.04)); border:1px solid rgba(0,229,255,.26);border-left:4px solid var(--c);border-radius:12px;padding:1.05rem 1.35rem;color:var(--c);font-family:'Exo 2',sans-serif;box-shadow:0 0 26px rgba(0,229,255,.05);}
@keyframes dp{0%,100%{box-shadow:0 0 26px rgba(255,59,92,.11)}50%{box-shadow:0 0 48px rgba(255,59,92,.28)}}

/* ── SECTION HEADER ───────────────────────── */
.shdr{font-family:'Share Tech Mono',monospace;font-size:.64rem;color:rgba(0,229,255,.46);
      text-transform:uppercase;letter-spacing:.24em;margin:1.5rem 0 .65rem 0;
      padding-bottom:.42rem;border-bottom:1px solid rgba(0,229,255,.09);position:relative;}
.shdr::after{content:'';position:absolute;bottom:-1px;left:0;width:52px;height:1px;
             background:var(--c);box-shadow:0 0 7px rgba(0,229,255,.5);}

/* ── PHASE BADGES ─────────────────────────── */
.pb{display:inline-block;font-family:'Share Tech Mono',monospace;font-size:.61rem;
    padding:.26rem .82rem;border-radius:20px;letter-spacing:.13em;text-transform:uppercase;
    margin-right:.52rem;font-weight:600;}
.pb1{background:rgba(0,229,255,.10);color:var(--c);border:1px solid var(--c);
     box-shadow:0 0 9px rgba(0,229,255,.22);text-shadow:0 0 6px rgba(0,229,255,.7);}
.pb2{background:rgba(147,51,234,.12);color:#c084fc;border:1px solid #9333ea;
     box-shadow:0 0 9px rgba(147,51,234,.26);text-shadow:0 0 6px rgba(147,51,234,.7);}
.pb3{background:rgba(0,255,136,.10);color:var(--g);border:1px solid var(--g);
     box-shadow:0 0 9px rgba(0,255,136,.22);text-shadow:0 0 6px rgba(0,255,136,.7);}

/* ── CONTENT BOX ──────────────────────────── */
.cb{background:rgba(2,8,20,.93);border:1px solid rgba(0,229,255,.10);
    border-radius:12px;padding:1.25rem;font-family:'Share Tech Mono',monospace;
    font-size:.79rem;line-height:1.9;color:#a8d4f0;white-space:pre-wrap;
    word-break:break-word;max-height:480px;overflow-y:auto;
    box-shadow:inset 0 0 38px rgba(0,0,0,.55);}
.cb::-webkit-scrollbar{width:4px;}
.cb::-webkit-scrollbar-thumb{background:linear-gradient(180deg,var(--c),var(--p));border-radius:2px;}

/* ── HEX DUMP ─────────────────────────────── */
.hx{background:rgba(0,5,15,.97);border:1px solid rgba(0,255,136,.16);
    border-radius:10px;padding:1.05rem 1.25rem;font-family:'Share Tech Mono',monospace;
    font-size:.74rem;color:#00ff88;line-height:1.8;overflow-x:auto;
    box-shadow:0 0 16px rgba(0,255,136,.04),inset 0 0 38px rgba(0,0,0,.8);
    text-shadow:0 0 5px rgba(0,255,136,.36);}

/* ── PPTX SLIDE CARD ──────────────────────── */
.sl-card{
  background:linear-gradient(135deg,rgba(10,20,50,.86),rgba(5,10,30,.86));
  border:1px solid rgba(147,51,234,.22);border-radius:14px;
  padding:1.35rem 1.55rem;margin-bottom:.9rem;position:relative;overflow:hidden;
  box-shadow:0 4px 28px rgba(147,51,234,.06);transition:border-color .2s,box-shadow .2s;
}
.sl-card:hover{border-color:rgba(147,51,234,.42);box-shadow:0 8px 38px rgba(147,51,234,.13);}
.sl-card::before{content:'';position:absolute;top:0;left:0;right:0;height:1px;
                  background:linear-gradient(90deg,transparent,rgba(147,51,234,.46),transparent);}
.sl-num{font-family:'Orbitron',monospace;font-size:.62rem;color:rgba(147,51,234,.68);
        letter-spacing:.20em;text-transform:uppercase;margin-bottom:.55rem;}
.sl-title{font-family:'Orbitron',monospace;font-size:1.02rem;font-weight:700;
          color:var(--c);margin-bottom:.75rem;line-height:1.3;}
.sl-body{font-family:'Exo 2',sans-serif;font-size:.86rem;color:#a8c8e8;line-height:1.8;}
.sl-bul{padding:.18rem 0;padding-left:1.15rem;position:relative;}
.sl-bul::before{content:'›';position:absolute;left:0;color:var(--p);font-weight:700;}
.sl-trow{padding:.18rem 0;padding-left:1.4rem;position:relative;
         color:#7c9fc0;font-family:'Share Tech Mono',monospace;font-size:.76rem;}
.sl-trow::before{content:'⊞';position:absolute;left:0;color:rgba(0,255,136,.5);}

/* ── FEATURE CARDS ────────────────────────── */
.fc{background:var(--gl);border:1px solid var(--bd);border-radius:15px;
    padding:1.4rem;margin-bottom:.7rem;backdrop-filter:blur(12px);
    transition:all .3s ease;position:relative;overflow:hidden;}
.fc:hover{border-color:rgba(0,229,255,.28);box-shadow:0 8px 38px rgba(0,229,255,.07);transform:translateY(-3px);}
.fc::before{content:'';position:absolute;top:0;left:0;right:0;height:1px;
             background:linear-gradient(90deg,transparent,rgba(0,229,255,.28),transparent);}
.fc-icon{font-size:1.75rem;margin-bottom:.55rem;}
.fc-title{font-family:'Orbitron',monospace;font-size:.76rem;font-weight:700;
          color:var(--c);letter-spacing:.10em;margin-bottom:.45rem;}
.fc-desc{font-family:'Exo 2',sans-serif;font-size:.81rem;color:rgba(160,200,240,.68);line-height:1.65;}

/* ── BLOCKED BOX ──────────────────────────── */
.blocked{background:rgba(255,59,92,.05);border:1px solid rgba(255,59,92,.28);
         border-radius:14px;padding:1.4rem 1.6rem;margin-top:.9rem;}

/* ── UPLOAD ZONE ──────────────────────────── */
[data-testid="stFileUploader"]{
  background:rgba(0,20,50,.36) !important;
  border:2px dashed rgba(0,229,255,.26) !important;
  border-radius:16px !important;backdrop-filter:blur(10px) !important;
  transition:all .3s ease !important;
  box-shadow:0 0 26px rgba(0,229,255,.03),inset 0 0 26px rgba(0,0,0,.26) !important;
}
[data-testid="stFileUploader"]:hover{
  border-color:rgba(0,229,255,.62) !important;
  box-shadow:0 0 38px rgba(0,229,255,.09),inset 0 0 26px rgba(0,229,255,.03) !important;
  background:rgba(0,30,70,.46) !important;
}
[data-testid="stFileUploader"] *{color:rgba(0,229,255,.76) !important;}

/* ── BUTTONS ──────────────────────────────── */
.stDownloadButton>button,
.stButton>button{
  background:linear-gradient(135deg,#002555,#07062a) !important;
  border:1px solid var(--c) !important;color:var(--c) !important;
  border-radius:10px !important;font-family:'Share Tech Mono',monospace !important;
  font-size:.78rem !important;letter-spacing:.10em !important;
  padding:.48rem 1.25rem !important;
  transition:all .22s ease !important;
  text-shadow:0 0 7px rgba(0,229,255,.52) !important;
  box-shadow:0 0 13px rgba(0,229,255,.12),inset 0 1px 0 rgba(0,229,255,.08) !important;
}
.stDownloadButton>button:hover,
.stButton>button:hover{
  background:linear-gradient(135deg,#003c7a,#160a5a) !important;
  box-shadow:0 0 26px rgba(0,229,255,.30),0 0 55px rgba(0,229,255,.10) !important;
  transform:translateY(-2px) !important;
  border-color:rgba(0,229,255,.86) !important;
}
/* Sidebar button — red */
[data-testid="stSidebar"] .stButton>button{
  background:linear-gradient(135deg,#2e0010,#180008) !important;
  border-color:rgba(255,59,92,.55) !important;color:var(--r) !important;
  text-shadow:0 0 7px rgba(255,59,92,.52) !important;
  box-shadow:0 0 13px rgba(255,59,92,.11) !important;
}
[data-testid="stSidebar"] .stButton>button:hover{
  box-shadow:0 0 26px rgba(255,59,92,.28) !important;
  border-color:var(--r) !important;transform:translateY(-2px) !important;
}

/* ── TABS ─────────────────────────────────── */
.stTabs [data-baseweb="tab-list"]{
  background:rgba(2,8,20,.82) !important;border-radius:12px !important;
  padding:5px !important;border:1px solid rgba(0,229,255,.09) !important;gap:2px !important;
}
.stTabs [data-baseweb="tab"]{
  color:var(--mu) !important;font-family:'Share Tech Mono',monospace !important;
  font-size:.73rem !important;border-radius:8px !important;padding:.38rem .95rem !important;
  letter-spacing:.05em !important;transition:all .2s !important;border:none !important;
}
.stTabs [data-baseweb="tab"]:hover{color:var(--c) !important;background:rgba(0,229,255,.05) !important;}
.stTabs [aria-selected="true"]{
  background:linear-gradient(135deg,rgba(0,229,255,.14),rgba(147,51,234,.09)) !important;
  color:var(--c) !important;box-shadow:0 0 11px rgba(0,229,255,.18) !important;
  border:1px solid rgba(0,229,255,.23) !important;text-shadow:0 0 5px rgba(0,229,255,.48) !important;
}

/* ── EXPANDERS ────────────────────────────── */
div[data-testid="stExpander"]{
  background:var(--gl2) !important;border:1px solid var(--bd) !important;
  border-radius:12px !important;backdrop-filter:blur(10px) !important;
}
div[data-testid="stExpander"] summary{
  color:var(--c) !important;font-family:'Share Tech Mono',monospace !important;font-size:.80rem !important;
}

/* ── INPUTS ───────────────────────────────── */
.stTextInput>div>div>input{
  background:rgba(0,10,30,.82) !important;
  border:1px solid rgba(0,229,255,.18) !important;
  border-radius:10px !important;color:var(--tx) !important;
  font-family:'Share Tech Mono',monospace !important;font-size:.82rem !important;
  caret-color:var(--c) !important;transition:border-color .2s,box-shadow .2s !important;
}
.stTextInput>div>div>input:focus{
  border-color:var(--c) !important;box-shadow:0 0 14px rgba(0,229,255,.14) !important;
}

/* ── DATAFRAME ────────────────────────────── */
[data-testid="stDataFrame"]{
  border:1px solid rgba(0,229,255,.09) !important;
  border-radius:12px !important;overflow:hidden !important;
}

/* ── SIDEBAR ELEMENTS ─────────────────────── */
.sb-logo{font-family:'Orbitron',monospace;font-size:1.15rem;font-weight:900;
         background:linear-gradient(135deg,var(--c),var(--p));
         -webkit-background-clip:text;-webkit-text-fill-color:transparent;
         background-clip:text;margin:0;letter-spacing:1px;}
.sb-tag{font-family:'Share Tech Mono',monospace;font-size:.58rem;
        color:rgba(0,229,255,.38);letter-spacing:.20em;text-transform:uppercase;margin:.18rem 0 1.4rem 0;}
.sb-ok{background:linear-gradient(135deg,rgba(0,255,136,.08),rgba(0,200,100,.04));
       border:1px solid rgba(0,255,136,.28);border-radius:10px;
       padding:.62rem .95rem;font-family:'Share Tech Mono',monospace;font-size:.66rem;
       color:var(--g);line-height:1.85;box-shadow:0 0 18px rgba(0,255,136,.05);}
.sb-warn{background:linear-gradient(135deg,rgba(255,214,10,.08),rgba(255,160,0,.04));
         border:1px solid rgba(255,214,10,.28);border-radius:10px;
         padding:.62rem .95rem;font-family:'Share Tech Mono',monospace;font-size:.66rem;
         color:var(--y);line-height:1.85;}
.lg{font-family:'Share Tech Mono',monospace;font-size:.64rem;color:rgba(74,107,138,.88);
    padding:.32rem 0;border-bottom:1px solid rgba(0,229,255,.055);}
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
#  CONSTANTS
# ══════════════════════════════════════════════════════════════
EXT_TO_MIME = {
    ".pdf" :[" application/pdf"],
    ".docx":["application/vnd.openxmlformats-officedocument.wordprocessingml.document","application/zip"],
    ".xlsx":["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","application/zip"],
    ".pptx":["application/vnd.openxmlformats-officedocument.presentationml.presentation","application/zip"],
    ".doc" :["application/msword"],
    ".xls" :["application/vnd.ms-excel"],
    ".ppt" :["application/vnd.ms-powerpoint"],
    ".txt" :["text/plain"],
    ".csv" :["text/plain","text/csv"],
    ".json":["application/json","text/plain"],
    ".xml" :["text/xml","application/xml"],
    ".png" :["image/png"],".jpg":["image/jpeg"],".jpeg":["image/jpeg"],
    ".gif" :["image/gif"],".bmp":["image/bmp","image/x-bmp"],".webp":["image/webp"],
    ".zip" :["application/zip"],
}
# strip accidental leading space in pdf entry
EXT_TO_MIME[".pdf"] = ["application/pdf"]

HIGH_RISK = {
    "application/x-dosexec"      :"Windows Executable (PE32/PE64)",
    "application/x-executable"   :"Linux ELF Executable",
    "application/x-sharedlib"    :"Shared Library (.so/.dll)",
    "text/x-shellscript"         :"Shell / Bash Script",
    "text/x-python"              :"Python Script",
    "application/x-msdos-program":"MS-DOS Program",
    "application/x-object"       :"Compiled Object Code",
}

MANUAL_SIGS = [
    (b"%PDF"             ,"application/pdf"         ,"PDF Document"),
    (b"\x89PNG\r\n\x1a\n","image/png"               ,"PNG Image"),
    (b"\xff\xd8\xff"     ,"image/jpeg"              ,"JPEG Image"),
    (b"GIF8"             ,"image/gif"               ,"GIF Image"),
    (b"BM"               ,"image/bmp"               ,"BMP Image"),
    (b"PK\x03\x04"       ,"application/zip"         ,"ZIP Archive"),
    (b"MZ"               ,"application/x-dosexec"   ,"Windows PE Executable"),
    (b"\x7fELF"          ,"application/x-executable","Linux ELF Executable"),
    (b"\x1f\x8b"         ,"application/gzip"        ,"Gzip Archive"),
    (b"ID3"              ,"audio/mpeg"              ,"MP3 Audio"),
    (b"<!DOCTYPE"        ,"text/html"               ,"HTML Document"),
    (b"<html"            ,"text/html"               ,"HTML Document"),
    (b"<?xml"            ,"text/xml"                ,"XML Document"),
    (b"#!"               ,"text/x-shellscript"      ,"Shell Script"),
    (b"RIFF"             ,"audio/wav"               ,"WAV Audio / RIFF"),
]


# ══════════════════════════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════════════════════════
def sha256_b(data:bytes)->str:
    return hashlib.sha256(data).hexdigest()

def fmt_bytes(n):
    for u in ["B","KB","MB","GB"]:
        if n<1024: return f"{n:.1f} {u}"
        n/=1024
    return f"{n:.1f} TB"

def verdict_cls(risk):
    return {"CRITICAL":"v-danger","WARNING":"v-warn","SAFE":"v-safe","INFO":"v-info"}.get(risk,"v-info")

def hex_dump(data:bytes,n:int=128)->str:
    lines=[]
    for i in range(0,min(len(data),n),16):
        c=data[i:i+16]
        lines.append(f"{i:04x}  {' '.join(f'{b:02x}' for b in c):<48}  {''.join(chr(b) if 32<=b<127 else '.' for b in c)}")
    return "\n".join(lines)

def dedup_cols(cols):
    """Fix duplicate / None / empty column names — prevents ValueError in st.dataframe."""
    seen={}; out=[]
    for c in cols:
        k = str(c).strip() if (c is not None and str(c).strip() not in ("","nan")) else "col"
        if k in seen:
            seen[k]+=1; out.append(f"{k}_{seen[k]}")
        else:
            seen[k]=0; out.append(k)
    return out

def safe_df(df):
    """Deduplicate cols, reset index, convert all to str — prevents Arrow/serialization errors."""
    import pandas as pd
    df = df.copy()
    df.columns = dedup_cols(list(df.columns))
    df = df.reset_index(drop=True)
    # Convert every column to string to avoid mixed-type Arrow issues
    for col in df.columns:
        df[col] = df[col].astype(str).replace({"nan":"","None":"","<NA>":""})
    return df


# ══════════════════════════════════════════════════════════════
#  DETECTION
# ══════════════════════════════════════════════════════════════
@st.cache_data(show_spinner=False, max_entries=30)
def detect_true_type(file_bytes:bytes, filename:str)->dict:
    result={"mime":"application/octet-stream","description":"Unknown binary",
            "confidence":"low","method":"fallback"}
    # 1. libmagic — most accurate
    try:
        import magic
        suffix = Path(filename).suffix or ".bin"
        fd, tmp_path = tempfile.mkstemp(suffix=suffix)
        try:
            os.write(fd, file_bytes)
            os.close(fd)
            result["mime"]        = magic.Magic(mime=True).from_file(tmp_path)
            result["description"] = magic.Magic(mime=False).from_file(tmp_path)
            result["confidence"]  = "high"
            result["method"]      = "libmagic"
        finally:
            try: os.unlink(tmp_path)
            except: pass
        return result
    except Exception:
        pass
    # 2. Manual header
    raw=file_bytes[:512]
    for sig,mime,desc in MANUAL_SIGS:
        if raw.startswith(sig):
            result.update({"mime":mime,"description":desc,"confidence":"medium","method":"manual_header"})
            return result
    # 3. Content sniff
    try:
        text=raw.decode("utf-8",errors="strict"); s=text.strip()
        if s.startswith(("{","[")): result.update({"mime":"application/json","description":"JSON Data","confidence":"medium","method":"sniff"})
        elif any(kw in text for kw in ["import ","def ","class "]): result.update({"mime":"text/x-python","description":"Python Script","confidence":"medium","method":"sniff"})
        else: result.update({"mime":"text/plain","description":"Plain Text","confidence":"low","method":"sniff"})
    except Exception: pass
    return result

def is_valid_office_zip(fb:bytes)->bool:
    try:
        with zipfile.ZipFile(io.BytesIO(fb)) as z: return "[Content_Types].xml" in z.namelist()
    except: return False

def get_office_sub(fb:bytes)->str:
    try:
        with zipfile.ZipFile(io.BytesIO(fb)) as z:
            n=z.namelist()
            if "word/document.xml"     in n: return "docx"
            if "xl/workbook.xml"       in n: return "xlsx"
            if "ppt/presentation.xml"  in n: return "pptx"
    except: pass
    return "zip"

def classify(true_mime:str, ext:str, fb:bytes)->dict:
    ext=ext.lower()
    expected=EXT_TO_MIME.get(ext,[])
    if true_mime in HIGH_RISK and ext not in [".exe",".elf",".sh",".py",".so",".bin",".dll"]:
        return {"risk":"CRITICAL","action":"BLOCK","label":"🔴 CRITICAL — SPOOFED FILE",
                "msg":f"Extension <b>{ext}</b> claims a safe document but magic bytes reveal: <b>{HIGH_RISK[true_mime]}</b>. File BLOCKED for your safety."}
    if true_mime=="application/zip" and ext in [".docx",".xlsx",".pptx"]:
        if is_valid_office_zip(fb):
            return {"risk":"SAFE","action":"RENDER","label":"✅ SAFE",
                    "msg":f"ZIP-based Office format <b>{ext.upper()}</b> verified — [Content_Types].xml confirmed."}
        return {"risk":"WARNING","action":"WARN","label":"🟡 WARNING",
                "msg":f"ZIP file with <b>{ext}</b> extension but missing Office XML structure. May be corrupt or spoofed."}
    if true_mime in expected:
        return {"risk":"SAFE","action":"RENDER","label":"✅ SAFE",
                "msg":f"File authentic. Magic bytes confirm <b>{true_mime}</b> matches <b>{ext}</b> extension."}
    if true_mime=="application/octet-stream":
        return {"risk":"WARNING","action":"WARN","label":"🟡 WARNING",
                "msg":"File type undetermined — possibly encrypted, corrupted, or unusual format."}
    if expected:
        return {"risk":"WARNING","action":"WARN","label":"🟡 WARNING — TYPE MISMATCH",
                "msg":f"Extension <b>{ext}</b> expects <code>{expected[0]}</code> but true type is <code>{true_mime}</code>."}
    return {"risk":"INFO","action":"RENDER","label":"ℹ️ UNKNOWN EXTENSION",
            "msg":f"Extension <b>{ext}</b> not in known list. File appears to be: <code>{true_mime}</code>"}


# ══════════════════════════════════════════════════════════════
#  RENDERERS
# ══════════════════════════════════════════════════════════════
def render_pdf(fb:bytes, max_pages:int=50)->dict:
    try:
        import fitz
        doc=fitz.open(stream=fb,filetype="pdf")
        total=doc.page_count; pages=[]
        for i in range(min(total,max_pages)):
            try: txt=doc[i].get_text().strip()
            except: txt=""
            pages.append({"num":i+1,"text":txt or "(Image-only page)","words":len(txt.split())})
        meta={k:v for k,v in (doc.metadata or {}).items() if v}
        doc.close()
        return {"ok":True,"type":"pdf","total_pages":total,"loaded_pages":len(pages),"pages":pages,"meta":meta}
    except ImportError: return {"ok":False,"error":"PyMuPDF not installed. Run: pip install PyMuPDF"}
    except Exception as e: return {"ok":False,"error":f"PDF error: {e}"}

def render_docx(fb:bytes)->dict:
    try:
        from docx import Document
        doc=Document(io.BytesIO(fb))
        paras=[{"style":p.style.name,"text":p.text.strip()} for p in doc.paragraphs if p.text.strip()]
        tables=[]
        for t in doc.tables:
            rows=[[c.text.strip() for c in r.cells] for r in t.rows]
            if rows: tables.append(rows)
        return {"ok":True,"type":"docx","paragraphs":paras,"tables":tables}
    except ImportError: return {"ok":False,"error":"python-docx not installed. Run: pip install python-docx"}
    except Exception as e: return {"ok":False,"error":f"DOCX error: {e}"}

def render_xlsx(fb:bytes)->dict:
    try:
        import pandas as pd
        sheets={}
        try: xl=pd.ExcelFile(io.BytesIO(fb),engine="openpyxl")
        except: xl=pd.ExcelFile(io.BytesIO(fb))
        for name in xl.sheet_names:
            try:
                df=xl.parse(name)
                sheets[str(name)]=safe_df(df)
            except Exception as se:
                sheets[str(name)+"__ERR"]=None
        return {"ok":True,"type":"xlsx","sheets":sheets}
    except ImportError: return {"ok":False,"error":"pandas/openpyxl not installed. Run: pip install pandas openpyxl"}
    except Exception as e: return {"ok":False,"error":f"XLSX error: {e}"}

def render_pptx(fb:bytes)->dict:
    try:
        from pptx import Presentation
        prs=Presentation(io.BytesIO(fb))
        slides=[]
        for i,slide in enumerate(prs.slides):
            title=""; items=[]; images=0; tables=0
            # Title
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title=slide.shapes.title.text.strip()
            except: pass
            # Shapes
            for shape in slide.shapes:
                try:
                    if shape.shape_type==13: images+=1; continue   # picture
                    if shape.has_table:
                        tables+=1
                        for row in shape.table.rows:
                            rt=" | ".join(c.text.strip() for c in row.cells if c.text.strip())
                            if rt: items.append({"level":2,"text":rt,"bold":False,"trow":True})
                        continue
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt=para.text.strip()
                            if not txt: continue
                            if txt==title and not items: continue
                            bold=False
                            try: bold=any(r.font.bold for r in para.runs if r.font.bold is True)
                            except: pass
                            items.append({"level":para.level,"text":txt,"bold":bold,"trow":False})
                except: pass
            if not title and items:
                bolds=[x for x in items if x.get("bold")]
                title=bolds[0]["text"] if bolds else items[0]["text"]
                items=[x for x in items if x["text"]!=title]
            slides.append({"num":i+1,"title":title or f"Slide {i+1}","items":items,"images":images,"tables":tables})
        meta={"slides":len(slides),"width_cm":round(prs.slide_width.cm,1),"height_cm":round(prs.slide_height.cm,1)}
        try:
            cp=prs.core_properties
            if cp.author:   meta["author"]  =cp.author
            if cp.title:    meta["title"]   =cp.title
            if cp.modified: meta["modified"]=str(cp.modified)[:19]
        except: pass
        return {"ok":True,"type":"pptx","slides":slides,"meta":meta}
    except ImportError: return {"ok":False,"error":"python-pptx not installed. Run: pip install python-pptx"}
    except Exception as e: return {"ok":False,"error":f"PPTX error: {e}\n{traceback.format_exc()}"}

def render_image(fb:bytes,filename:str)->dict:
    try:
        from PIL import Image
        img=Image.open(io.BytesIO(fb)); img.verify()
        img=Image.open(io.BytesIO(fb))
        return {"ok":True,"type":"image","format":img.format or "?","mode":img.mode,"size":img.size,"raw":fb}
    except ImportError: return {"ok":False,"error":"Pillow not installed. Run: pip install Pillow"}
    except Exception as e: return {"ok":False,"error":f"Image error: {e}"}

def render_text(fb:bytes,mime:str,filename:str)->dict:
    try:
        text=fb.decode("utf-8",errors="replace"); ext=Path(filename).suffix.lower()
        if mime=="application/json" or ext==".json":
            try: return {"ok":True,"type":"json","raw":text,"parsed":json.loads(text)}
            except: pass
        if ext==".csv" or mime in ("text/csv",):
            try:
                import pandas as pd
                df=pd.read_csv(io.StringIO(text))
                return {"ok":True,"type":"csv","raw":text,"df":safe_df(df)}
            except: pass
        return {"ok":True,"type":"text","raw":text,"lines":text.splitlines()}
    except Exception as e: return {"ok":False,"error":f"Text error: {e}"}

def render_zip(fb:bytes)->dict:
    try:
        with zipfile.ZipFile(io.BytesIO(fb)) as z:
            entries=[{"name":i.filename,"size":i.file_size,"compressed":i.compress_size,
                      "is_dir":i.filename.endswith("/")} for i in z.infolist()]
        return {"ok":True,"type":"zip","entries":entries}
    except Exception as e: return {"ok":False,"error":f"ZIP error: {e}"}

def smart_render(fb:bytes,true_mime:str,filename:str,max_pages:int=50)->dict:
    ext=Path(filename).suffix.lower()
    if true_mime=="application/pdf":                                          return render_pdf(fb,max_pages)
    if true_mime in["application/vnd.openxmlformats-officedocument.wordprocessingml.document"] or (true_mime=="application/zip" and ext==".docx"): return render_docx(fb)
    if true_mime in["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","application/vnd.ms-excel"] or (true_mime=="application/zip" and ext in[".xlsx",".xls"]): return render_xlsx(fb)
    if true_mime in["application/vnd.openxmlformats-officedocument.presentationml.presentation","application/vnd.ms-powerpoint"] or (true_mime=="application/zip" and ext in[".pptx",".ppt"]): return render_pptx(fb)
    if true_mime=="application/zip":
        sub=get_office_sub(fb)
        if sub=="docx": return render_docx(fb)
        if sub=="xlsx": return render_xlsx(fb)
        if sub=="pptx": return render_pptx(fb)
        return render_zip(fb)
    if true_mime in["image/png","image/jpeg","image/gif","image/bmp","image/webp","image/tiff"]: return render_image(fb,filename)
    if true_mime in["text/plain","text/csv","application/json","text/xml","application/xml","text/html","text/x-python","text/x-shellscript"]: return render_text(fb,true_mime,filename)
    if true_mime in["application/gzip","application/x-gzip"]: return{"ok":False,"error":"Gzip file — please decompress first."}
    return{"ok":False,"error":f"No renderer for: {true_mime}"}


# ══════════════════════════════════════════════════════════════
#  SESSION STATE
# ══════════════════════════════════════════════════════════════
if "history" not in st.session_state: st.session_state.history=[]
if "last_file_id" not in st.session_state: st.session_state.last_file_id=None


# ══════════════════════════════════════════════════════════════
#  SIDEBAR
# ══════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown('<p class="sb-logo">🛡️ DocShield</p>', unsafe_allow_html=True)
    st.markdown('<p class="sb-tag">Universal Document Reader  v4.0</p>', unsafe_allow_html=True)

    st.markdown('<p class="shdr">⚙️  Configuration</p>', unsafe_allow_html=True)
    max_pages    = st.slider("📄 Max PDF pages", 1, 200, 50)
    show_hex     = st.toggle("🔬 Show binary hex dump", False)
    force_render = st.toggle("⚡ Force render (override block)", False)

    st.markdown('<p class="shdr">🔮  Detection Engine</p>', unsafe_allow_html=True)
    try:
        import magic
        st.markdown('<div class="sb-ok">● libmagic · ACTIVE<br>✓ Binary Signature DB<br>✓ 1000+ Signatures<br>✓ High Confidence</div>', unsafe_allow_html=True)
    except ImportError:
        st.markdown('<div class="sb-warn">◐ libmagic · FALLBACK<br>⚠ Manual Header Scan<br>⚠ Medium Confidence</div>', unsafe_allow_html=True)

    st.markdown('<p class="shdr">📋  Scan History</p>', unsafe_allow_html=True)
    if st.session_state.history:
        for h in reversed(st.session_state.history[-10:]):
            icon="🔴" if h["risk"]=="CRITICAL" else "🟡" if h["risk"]=="WARNING" else "✅"
            st.markdown(f'<div class="lg">{icon} {h["name"][:20]}<br><span style="color:rgba(74,107,138,.45);font-size:.56rem">{h["time"]} · {h["mime"][:22]}</span></div>', unsafe_allow_html=True)
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🗑  Clear History"): st.session_state.history=[]; st.rerun()
    else:
        st.markdown('<p style="font-family:\'Share Tech Mono\',monospace;font-size:.66rem;color:rgba(74,107,138,.48)">No files scanned yet.</p>', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
st.markdown('<div class="hero"><p class="hero-title">🛡️ DocShield Reader</p><p class="hero-sub">Magic-Based File Integrity · PDF · DOCX · XLSX · PPTX · Images · JSON · ZIP · Zero Trust</p></div>', unsafe_allow_html=True)

# Stats
h=st.session_state.history
c1,c2,c3,c4=st.columns(4)
c1.markdown(f'<div class="sc sc-c"><div class="sc-lbl">🗂 Files Scanned</div><div class="sc-val" style="color:#00e5ff">{len(h)}</div></div>', unsafe_allow_html=True)
c2.markdown(f'<div class="sc sc-r"><div class="sc-lbl">⛔ Threats Blocked</div><div class="sc-val" style="color:#ff3b5c">{sum(1 for x in h if x["risk"]=="CRITICAL")}</div></div>', unsafe_allow_html=True)
c3.markdown(f'<div class="sc sc-y"><div class="sc-lbl">⚠️ Warnings</div><div class="sc-val" style="color:#ffd60a">{sum(1 for x in h if x["risk"]=="WARNING")}</div></div>', unsafe_allow_html=True)
c4.markdown(f'<div class="sc sc-g"><div class="sc-lbl">✅ Safe Files</div><div class="sc-val" style="color:#00ff88">{sum(1 for x in h if x["risk"]=="SAFE")}</div></div>', unsafe_allow_html=True)

# Upload
st.markdown('<p class="shdr">📁  Drop Any File for Analysis</p>', unsafe_allow_html=True)
uploaded=st.file_uploader("Upload", type=None, label_visibility="collapsed",
    help="PDF · DOCX · XLSX · PPTX · PNG · JPEG · JSON · CSV · TXT · ZIP — and detects disguised EXE/ELF/scripts")

if uploaded:
    # ── FIX: read bytes robustly, catch network/axios read errors ──
    try:
        file_bytes = uploaded.read()
    except Exception as read_err:
        st.error(f"❌ File read error (network/upload issue): {read_err}")
        st.info("💡 Try re-uploading the file. If error persists, check your network connection.")
        st.stop()

    if not file_bytes or len(file_bytes)==0:
        st.error("❌ Empty file (0 bytes). Please upload a valid file.")
        st.stop()

    filename = uploaded.name
    ext      = Path(filename).suffix
    fsize    = len(file_bytes)

    # ── FIX: cache hash to avoid repeated recomputation on widget interactions ──
    fhash = sha256_b(file_bytes)
    if st.session_state.last_file_id != fhash:
        st.session_state.last_file_id = fhash

    st.markdown("---")

    # ════════════════════════════════════════
    #  PHASE 1 — Magic Detection
    # ════════════════════════════════════════
    st.markdown('<span class="pb pb1">Phase 1</span><strong style="color:#e0f0ff;font-family:\'Exo 2\';font-size:.93rem"> Magic Byte Inspection</strong>', unsafe_allow_html=True)
    with st.spinner("🔍 Reading binary signature…"):
        detection = detect_true_type(file_bytes, filename)

    d1,d2,d3=st.columns(3)
    d1.markdown(f'<div class="ic"><div class="ic-lbl">True MIME Type</div><div class="ic-val" style="color:#00e5ff">{detection["mime"]}</div></div>', unsafe_allow_html=True)
    d2.markdown(f'<div class="ic"><div class="ic-lbl">File Description</div><div class="ic-val">{detection["description"]}</div></div>', unsafe_allow_html=True)
    d3.markdown(f'<div class="ic"><div class="ic-lbl">Method · Confidence</div><div class="ic-val"><span style="color:#00ff88">{detection["method"]}</span> · {detection["confidence"]}</div></div>', unsafe_allow_html=True)

    e1,e2,e3=st.columns(3)
    e1.markdown(f'<div class="ic"><div class="ic-lbl">Filename</div><div class="ic-val">{filename}</div></div>', unsafe_allow_html=True)
    e2.markdown(f'<div class="ic"><div class="ic-lbl">Claimed Extension</div><div class="ic-val" style="color:#ffd60a">{ext or "(none)"}</div></div>', unsafe_allow_html=True)
    e3.markdown(f'<div class="ic"><div class="ic-lbl">File Size</div><div class="ic-val">{fmt_bytes(fsize)}</div></div>', unsafe_allow_html=True)

    if show_hex:
        st.markdown('<p class="shdr">🔬  Raw Binary Header — First 128 Bytes</p>', unsafe_allow_html=True)
        st.markdown(f'<div class="hx">{hex_dump(file_bytes[:128])}</div>', unsafe_allow_html=True)

    # ════════════════════════════════════════
    #  PHASE 2 — Threat Classification
    # ════════════════════════════════════════
    st.markdown('<br><span class="pb pb2">Phase 2</span><strong style="color:#e0f0ff;font-family:\'Exo 2\';font-size:.93rem"> Threat Classification</strong>', unsafe_allow_html=True)
    verdict = classify(detection["mime"], ext, file_bytes)
    st.markdown(f'<div class="{verdict_cls(verdict["risk"])}" style="margin-top:.45rem"><strong style="font-size:.98rem;letter-spacing:.05em">{verdict["label"]}</strong><br><br>{verdict["msg"]}</div>', unsafe_allow_html=True)
    st.markdown(f'<div style="font-family:\'Share Tech Mono\',monospace;font-size:.62rem;color:rgba(74,107,138,.65);margin-top:.45rem">SHA-256: {fhash}</div>', unsafe_allow_html=True)

    if not any(x["hash"]==fhash for x in st.session_state.history):
        st.session_state.history.append({"name":filename,"risk":verdict["risk"],"mime":detection["mime"],"hash":fhash,"time":datetime.now().strftime("%H:%M:%S")})

    # ════════════════════════════════════════
    #  PHASE 3 — Render
    # ════════════════════════════════════════
    if verdict["action"]=="BLOCK" and not force_render:
        st.markdown('<div class="blocked"><strong style="color:#ff3b5c;font-family:\'Orbitron\',monospace;font-size:.88rem;letter-spacing:.10em">⛔ FILE QUARANTINED</strong><br><br><span style="color:#94a3b8;font-size:.85rem">Rendering disabled. Enable <strong style="color:#ffd60a">Force render</strong> in sidebar if this is a safe test file.</span></div>', unsafe_allow_html=True)
    else:
        if verdict["action"]=="WARN":
            st.warning("⚠️ File has anomalies — rendering with caution. Verify the source before trusting content.")

        st.markdown('<br><span class="pb pb3">Phase 3</span><strong style="color:#e0f0ff;font-family:\'Exo 2\';font-size:.93rem"> Content Rendering</strong>', unsafe_allow_html=True)

        with st.spinner("📄 Rendering…"):
            result = smart_render(file_bytes, detection["mime"], filename, max_pages)

        if not result.get("ok"):
            st.markdown(f'<div class="v-warn" style="margin-top:.45rem"><strong>⚠️ Render Error</strong><br><br>{result.get("error","Unknown error")}</div>', unsafe_allow_html=True)
        else:
            rtype=result.get("type","")

            # ── PDF ──────────────────────────────────────────────
            if rtype=="pdf":
                total_p=result["total_pages"]; loaded=result["loaded_pages"]; meta=result.get("meta",{})
                st.markdown(f'<div class="v-info" style="margin-top:.45rem">📄 PDF · <b>{total_p}</b> total pages · <b>{loaded}</b> extracted (max={max_pages})</div>', unsafe_allow_html=True)
                if meta:
                    with st.expander("📋 Document Metadata"):
                        mc1,mc2=st.columns(2)
                        for i,(k,v) in enumerate(meta.items()):
                            (mc1 if i%2==0 else mc2).markdown(f'<div class="ic"><div class="ic-lbl">{k}</div><div class="ic-val" style="font-size:.78rem">{v}</div></div>', unsafe_allow_html=True)
                pages=result["pages"]
                if pages:
                    st.markdown('<p class="shdr">🔍  Search All Pages</p>', unsafe_allow_html=True)
                    sq=st.text_input("Search pages","",placeholder="Keyword across all extracted pages…",label_visibility="collapsed",key="pdf_sq")
                    if sq.strip():
                        hits=[p for p in pages if sq.lower() in p["text"].lower()]
                        st.markdown(f'<div class="v-info" style="padding:.55rem 1rem">Found <b>{len(hits)}</b> page(s) containing "<b>{sq}</b>"</div>', unsafe_allow_html=True)
                        for p in hits[:10]:
                            with st.expander(f"📄 Page {p['num']}  ({p['words']} words)"):
                                st.markdown(f'<div class="cb">{p["text"][:4000]}</div>', unsafe_allow_html=True)
                    else:
                        tab_n=min(len(pages),30)
                        tlbls=[f"P{p['num']}" for p in pages[:tab_n]]
                        if len(pages)>30: tlbls.append(f"+{len(pages)-30} more")
                        ptabs=st.tabs(tlbls)
                        for i,page in enumerate(pages[:30]):
                            with ptabs[i]:
                                st.markdown(f'<div style="font-family:\'Share Tech Mono\',monospace;font-size:.66rem;color:#4a6b8a;margin-bottom:.4rem">Page {page["num"]} · {page["words"]} words</div>', unsafe_allow_html=True)
                                st.markdown(f'<div class="cb">{page["text"]}</div>', unsafe_allow_html=True)
                        if len(pages)>30:
                            with ptabs[30]: st.info(f"Pages 31–{loaded} extracted. Use search above.")
                ft="\n\n".join(f"=== PAGE {p['num']} ===\n{p['text']}" for p in pages)
                st.download_button("⬇️  Download Extracted Text (.txt)", ft.encode("utf-8"), file_name=f"{Path(filename).stem}_text.txt", mime="text/plain")

            # ── DOCX ─────────────────────────────────────────────
            elif rtype=="docx":
                paras=result["paragraphs"]; tables=result["tables"]
                st.markdown(f'<div class="v-info" style="margin-top:.45rem">📝 Word Document · <b>{len(paras)}</b> paragraphs · <b>{len(tables)}</b> tables</div>', unsafe_allow_html=True)
                tc,tt=st.tabs(["📝 Content","📊 Tables"])
                with tc:
                    out=""
                    for p in paras:
                        sty,txt=p["style"],p["text"]
                        if "Heading 1" in sty: st.markdown(f'<div style="font-family:\'Orbitron\',monospace;font-size:.98rem;color:#00e5ff;padding:.75rem 0 .18rem;border-bottom:1px solid rgba(0,229,255,.14);margin-bottom:.35rem">{txt}</div>', unsafe_allow_html=True)
                        elif "Heading 2" in sty: st.markdown(f'<div style="font-family:\'Orbitron\',monospace;font-size:.82rem;color:#c084fc;padding:.45rem 0 .18rem">{txt}</div>', unsafe_allow_html=True)
                        elif "Heading" in sty: st.markdown(f'<div style="font-size:.80rem;color:#64748b;font-weight:600;padding:.26rem 0">{txt}</div>', unsafe_allow_html=True)
                        else: st.markdown(f'<div style="font-family:\'Exo 2\',sans-serif;font-size:.82rem;color:#a8c8e8;padding:.10rem 0 .10rem .45rem;line-height:1.7">{txt}</div>', unsafe_allow_html=True)
                        out+=f"{txt}\n"
                    st.download_button("⬇️  Download as Text", out.encode(), file_name=f"{Path(filename).stem}.txt", mime="text/plain")
                with tt:
                    if tables:
                        for ti,tbl in enumerate(tables):
                            st.markdown(f'<div style="font-family:\'Orbitron\',monospace;font-size:.72rem;color:#00e5ff;margin:.7rem 0 .35rem">TABLE {ti+1}</div>', unsafe_allow_html=True)
                            try:
                                import pandas as pd
                                if len(tbl)>1:
                                    df=pd.DataFrame(tbl[1:],columns=dedup_cols([str(c) for c in tbl[0]]))
                                    st.dataframe(safe_df(df),use_container_width=True)
                                else: st.write(tbl)
                            except Exception as te: st.write(tbl); st.caption(f"Table display err: {te}")
                    else: st.markdown('<p style="color:#4a6b8a;font-family:\'Share Tech Mono\',monospace;font-size:.78rem">No tables found.</p>', unsafe_allow_html=True)

            # ── XLSX ─────────────────────────────────────────────
            elif rtype=="xlsx":
                sheets=result["sheets"]
                valid={k:v for k,v in sheets.items() if v is not None}
                failed=[k for k,v in sheets.items() if v is None]
                st.markdown(f'<div class="v-info" style="margin-top:.45rem">📊 Excel · <b>{len(sheets)}</b> sheet(s)</div>', unsafe_allow_html=True)
                if failed: st.warning(f"⚠️ Failed to parse: {', '.join(failed)}")
                if valid:
                    if len(valid)==1:
                        name,df=next(iter(valid.items()))
                        st.markdown(f'<div style="font-family:\'Share Tech Mono\',monospace;font-size:.70rem;color:#4a6b8a;margin-bottom:.45rem">Sheet: <span style="color:#00e5ff">{name}</span> — {len(df)} rows × {len(df.columns)} cols</div>', unsafe_allow_html=True)
                        st.dataframe(df,use_container_width=True,height=420)
                    else:
                        tabs=st.tabs([f"📋 {n}" for n in valid.keys()])
                        for i,(name,df) in enumerate(valid.items()):
                            with tabs[i]:
                                st.markdown(f'<div style="font-size:.70rem;color:#4a6b8a;margin-bottom:.38rem">{len(df)} rows × {len(df.columns)} cols</div>', unsafe_allow_html=True)
                                st.dataframe(df,use_container_width=True,height=400)

            # ── PPTX ─────────────────────────────────────────────
            elif rtype=="pptx":
                slides=result["slides"]; meta=result.get("meta",{})
                st.markdown(f'<div class="v-info" style="margin-top:.45rem">📊 PowerPoint · <b>{meta.get("slides",len(slides))}</b> slides · {meta.get("width_cm","?")} × {meta.get("height_cm","?")} cm</div>', unsafe_allow_html=True)
                if len(meta)>3:
                    with st.expander("📋 Presentation Metadata"):
                        pm1,pm2=st.columns(2)
                        extra=[(k,v) for k,v in meta.items() if k not in["slides","width_cm","height_cm"]]
                        for i,(k,v) in enumerate(extra):
                            (pm1 if i%2==0 else pm2).markdown(f'<div class="ic"><div class="ic-lbl">{k}</div><div class="ic-val" style="font-size:.78rem">{v}</div></div>', unsafe_allow_html=True)
                st.markdown('<p class="shdr">🔍  Search Slides</p>', unsafe_allow_html=True)
                sq2=st.text_input("Search slides","",placeholder="Search text across all slides…",label_visibility="collapsed",key="pptx_sq")
                display_slides=slides
                if sq2.strip():
                    display_slides=[s for s in slides if sq2.lower() in s["title"].lower() or any(sq2.lower() in it["text"].lower() for it in s["items"])]
                    st.markdown(f'<div class="v-info" style="padding:.55rem 1rem;margin-bottom:.6rem">Found <b>{len(display_slides)}</b> slide(s) matching "<b>{sq2}</b>"</div>', unsafe_allow_html=True)
                st.markdown(f'<p class="shdr">🗂  Slides ({len(display_slides)} shown)</p>', unsafe_allow_html=True)
                for slide in display_slides:
                    extras = ""
                    if slide["images"]: extras += f"  🖼 {slide['images']} img"
                    if slide["tables"]: extras += f"  📊 {slide['tables']} tbl"
                    body_html = ""
                    for item in slide["items"]:
                        ind = item["level"] * 18
                        b1 = "<b>" if item.get("bold") else ""
                        b2 = "</b>" if item.get("bold") else ""
                        if item.get("trow"):
                            body_html += f'<div class="sl-trow" style="padding-left:{ind+22}px">{item["text"]}</div>'
                        else:
                            bul = "&#9670;" if item["level"] == 0 else ("&#9658;" if item["level"] == 1 else "&ndash;")
                            body_html += f'<div class="sl-bul" style="padding-left:{ind+18}px">{bul} {b1}{item["text"]}{b2}</div>'
                    no_content = '<span style="color:#334155;font-style:italic">No text content.</span>'
                    st.markdown(
                        f'<div class="sl-card">'
                        f'<div class="sl-num">SLIDE {slide["num"]}{extras}</div>'
                        f'<div class="sl-title">{slide["title"]}</div>'
                        f'<div class="sl-body">{body_html if body_html else no_content}</div>'
                        f'</div>',
                        unsafe_allow_html=True
                    )
                pt_lines = []
                for s in slides:
                    pt_lines.append(f"{'='*48}")
                    pt_lines.append(f"SLIDE {s['num']}: {s['title']}")
                    pt_lines.append(f"{'='*48}")
                    for it in s["items"]:
                        pt_lines.append(f"  {'  '*it['level']}* {it['text']}")
                pt = "\n".join(pt_lines)
                st.download_button("⬇️  Download Slides as Text", pt.encode("utf-8"), file_name=f"{Path(filename).stem}_slides.txt", mime="text/plain")

            # ── IMAGE ─────────────────────────────────────────────
            elif rtype=="image":
                ic1,ic2=st.columns([2,1])
                with ic1: st.image(result["raw"],use_container_width=True,caption=filename)
                with ic2:
                    for lbl,val in [("Format",result["format"]),("Color Mode",result["mode"]),("Width",f"{result['size'][0]} px"),("Height",f"{result['size'][1]} px"),("File Size",fmt_bytes(fsize))]:
                        st.markdown(f'<div class="ic"><div class="ic-lbl">{lbl}</div><div class="ic-val">{val}</div></div>', unsafe_allow_html=True)

            # ── JSON ──────────────────────────────────────────────
            elif rtype=="json":
                st.markdown('<div class="v-info" style="margin-top:.45rem">✅ Valid JSON Document</div>', unsafe_allow_html=True)
                st.json(result["parsed"])

            # ── CSV ───────────────────────────────────────────────
            elif rtype=="csv":
                df=result.get("df")
                if df is not None:
                    st.markdown(f'<div class="v-info" style="margin-top:.45rem">📊 CSV · <b>{len(df)}</b> rows · <b>{len(df.columns)}</b> columns</div>', unsafe_allow_html=True)
                    st.dataframe(df,use_container_width=True,height=420)
                else:
                    st.markdown(f'<div class="cb">{result["raw"][:6000]}</div>', unsafe_allow_html=True)

            # ── TEXT / SCRIPT ─────────────────────────────────────
            elif rtype=="text":
                lines=result["lines"]
                st.markdown(f'<div class="v-info" style="margin-top:.45rem">📄 Text · <b>{len(lines)}</b> lines · <b>{len(result["raw"]):,}</b> chars</div>', unsafe_allow_html=True)
                if detection["mime"] in["text/x-python","text/x-shellscript"]:
                    st.warning("⚠️ Script file — read-only display. NOT executing.")
                lang="python" if "python" in detection["mime"] else "bash" if "shell" in detection["mime"] else "text"
                st.code(result["raw"][:15000],language=lang)
                if len(result["raw"])>15000: st.info(f"Showing first 15,000 of {len(result['raw']):,} chars.")

            # ── ZIP ───────────────────────────────────────────────
            elif rtype=="zip":
                import pandas as pd
                entries=result["entries"]
                files=[e for e in entries if not e["is_dir"]]
                dirs=[e for e in entries if e["is_dir"]]
                st.markdown(f'<div class="v-info" style="margin-top:.45rem">🗜 ZIP Archive · <b>{len(files)}</b> files · <b>{len(dirs)}</b> folders</div>', unsafe_allow_html=True)
                if files:
                    df_z=pd.DataFrame([{"Filename":e["name"],"Size":fmt_bytes(e["size"]),"Compressed":fmt_bytes(e["compressed"])} for e in files])
                    st.dataframe(df_z,use_container_width=True,height=400)

else:
    # ── EMPTY STATE ──────────────────────────────────────────────
    st.markdown('<div style="text-align:center;padding:2.2rem 1rem 1.3rem"><div style="font-size:3.3rem;margin-bottom:.7rem;filter:drop-shadow(0 0 20px rgba(0,229,255,.5))">🛡️</div><p style="font-family:\'Orbitron\',monospace;font-size:1.02rem;font-weight:700;background:linear-gradient(90deg,#00e5ff,#9333ea,#ec4899);-webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text;letter-spacing:2px;margin-bottom:.4rem">DROP ANY FILE TO ANALYZE</p><p style="font-family:\'Share Tech Mono\',monospace;font-size:.68rem;color:rgba(0,229,255,.36);letter-spacing:.20em;text-transform:uppercase;max-width:540px;margin:0 auto;line-height:2.2">PDF · DOCX · XLSX · PPTX · PNG · JPEG · JSON · CSV · TXT · ZIP<br>Detects → EXE · ELF · Shell Scripts · Python · Malware</p></div>', unsafe_allow_html=True)
    fa,fb,fc=st.columns(3)
    fa.markdown('<div class="fc"><div class="fc-icon">🔬</div><div class="fc-title">Magic Detection</div><div class="fc-desc">Reads raw binary header. Identifies 1000+ types. Cannot be fooled by renaming malware.exe to report.pdf.</div></div>', unsafe_allow_html=True)
    fb.markdown('<div class="fc"><div class="fc-icon">⛔</div><div class="fc-title">Threat Blocking</div><div class="fc-desc">EXE, ELF, Shell scripts, Python files disguised as documents are blocked instantly with SHA-256 logging.</div></div>', unsafe_allow_html=True)
    fc.markdown('<div class="fc"><div class="fc-icon">📊</div><div class="fc-title">PPTX Support</div><div class="fc-desc">Full PowerPoint renderer — every slide title, bullets, table rows, image counts in beautiful slide cards.</div></div>', unsafe_allow_html=True)
    fd,fe,ff=st.columns(3)
    fd.markdown('<div class="fc"><div class="fc-icon">📄</div><div class="fc-title">PDF · 200 Pages</div><div class="fc-desc">Extract up to 200 pages with per-page tabs, full-text keyword search, highlighted results, one-click download.</div></div>', unsafe_allow_html=True)
    fe.markdown('<div class="fc"><div class="fc-icon">🔧</div><div class="fc-title">Zero ValueError</div><div class="fc-desc">Duplicate columns auto-fixed. Multi-sheet Excel. Arrow serialization errors prevented. No crashes.</div></div>', unsafe_allow_html=True)
    ff.markdown('<div class="fc"><div class="fc-icon">⚡</div><div class="fc-title">Zero Trust Mode</div><div class="fc-desc">Extension is only a hint. Binary magic bytes are the only truth. File name means nothing — bytes mean everything.</div></div>', unsafe_allow_html=True)
